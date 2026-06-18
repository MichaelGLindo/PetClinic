from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
TEAL = RGBColor(20, 184, 166)
VIOLET = RGBColor(139, 92, 246)
DARK_BG = RGBColor(10, 15, 30)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)
DARK_CARD = RGBColor(26, 32, 53)
GREEN = RGBColor(34, 197, 94)
RED = RGBColor(239, 68, 68)
YELLOW = RGBColor(234, 179, 8)

def set_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox

def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    table = slide.shapes.add_table(len(rows)+1, len(headers), left, top, width, height).table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r+1, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_CARD if r % 2 == 0 else RGBColor(20, 28, 48)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = LIGHT_GRAY
    return table

# ===================== SLIDE 1: PORTADA =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5), "PetClinic", 54, TEAL, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1), "Sistema de Gestion Veterinaria", 28, WHITE, False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8), "Arquitectura, Backend, Frontend, Seguridad y Automatizacion", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6), "Integrante: Michael Galindo  |  Junio 2026", 16, VIOLET, False, PP_ALIGN.CENTER)

# ===================== SLIDE 2: QUE ES PETCLINIC =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "1. Que es PetClinic?", 32, TEAL, True)

add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2), DARK_CARD)
add_textbox(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.5), "Descripcion", 20, VIOLET, True)
add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5), Inches(4), [
    "- Sistema completo de gestion para veterinarias",
    "- Administra: Dueños, Mascotas, Turnos e Historial Clinico",
    "- Backend API + Frontend Web + Automatizacion de pruebas",
    "- App movil: pendiente de implementar",
], 15, LIGHT_GRAY)

add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.2), DARK_CARD)
add_textbox(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.5), "Estructura del Proyecto", 20, VIOLET, True)
add_bullet_list(slide, Inches(7.1), Inches(2.3), Inches(5.2), Inches(4), [
    "- backend/  →  API REST con Python FastAPI",
    "- frontend-web/  →  App web con React 19",
    "- PetClinicAutomation/  →  Serenity BDD + Cucumber",
    "- DOCUMENTACION/  →  Diagramas y docs",
    "- app-movil/  →  Vacia (pendiente)",
], 15, LIGHT_GRAY)

# ===================== SLIDE 3: STACK TECNOLOGICO =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "2. Stack Tecnologico", 32, TEAL, True)

headers = ["Componente", "Tecnologia", "Estado"]
rows = [
    ["Backend", "Python FastAPI + SQLAlchemy", "Activo"],
    ["Frontend", "React 19 + React Router v7", "Activo"],
    ["Base de Datos", "MySQL (veterinaria)", "Activo"],
    ["Autenticacion", "JWT (HS256, 60min)", "Activo"],
    ["Automatizacion", "Serenity BDD + Cucumber", "Activo"],
    ["App Movil", "-", "Pendiente"],
]
add_table(slide, Inches(1), Inches(1.6), Inches(11), Inches(4.5), headers, rows)

# ===================== SLIDE 4: MODELO DE DOMINIO =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "3. Modelo de Dominio (5 Entidades)", 32, TEAL, True)

headers = ["Entidad", "Tabla BD", "Campos Principales"]
rows = [
    ["Dueno", "duenos", "cedula (PK), nombre, telefono"],
    ["Mascota", "mascotas", "id (PK), nombre, especie, edad, dueno_cedula (FK)"],
    ["Turno", "turnos", "id (PK), fecha, motivo, mascota_id (FK)"],
    ["Historial", "historial_clinico", "id (PK), fecha, descripcion, diagnostico, mascota_id (FK)"],
    ["Usuario", "usuarios", "id (PK), username, password_hash, rol, dueno_cedula (FK)"],
]
add_table(slide, Inches(1), Inches(1.5), Inches(11), Inches(3.5), headers, rows)

add_textbox(slide, Inches(0.8), Inches(5.3), Inches(5), Inches(0.5), "Relaciones:", 18, VIOLET, True)
add_bullet_list(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1.5), [
    "- Dueno (1) → (*) Mascota  |  Mascota (1) → (*) Turno  |  Mascota (1) → (*) Historial",
    "- Dueno (1) → (1) Usuario  |  ON DELETE CASCADE en todas las foreign keys",
], 14, LIGHT_GRAY)

# ===================== SLIDE 5: BACKEND FASTAPI =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "4. Backend - Python FastAPI", 32, TEAL, True)

add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.3), DARK_CARD)
add_textbox(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.5), "Arquitectura en Capas", 20, VIOLET, True)
add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5), Inches(4.2), [
    "- Router Layer: endpoints HTTP, respuestas JSON",
    "- CRUD Layer: operaciones de base de datos",
    "- Model Layer: entidades con SQLAlchemy ORM",
    "- Schema Layer: validacion con Pydantic (DTOs)",
    "- Security: JWT + bcrypt + roles",
], 15, LIGHT_GRAY)

add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), DARK_CARD)
add_textbox(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.5), "Archivos Principales", 20, VIOLET, True)
add_bullet_list(slide, Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2), [
    "- main.py →入口, CORS, routers, crea tablas",
    "- database.py → conexion MySQL, get_db()",
    "- models.py → modelos SQLAlchemy",
    "- schemas.py → schemas Pydantic",
    "- crud.py → operaciones CRUD",
    "- security.py → hash bcrypt, JWT",
    "- roles.py → require_role() autorizacion",
], 15, LIGHT_GRAY)

# ===================== SLIDE 6: ENDPOINTS API =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "5. Endpoints API REST", 32, TEAL, True)

headers = ["Metodo", "Ruta", "Descripcion"]
rows = [
    ["POST", "/auth/login", "Login, retorna token JWT"],
    ["POST", "/auth/register-client", "Registrar cliente + dueo"],
    ["GET/POST/PUT/DELETE", "/api/duenos", "CRUD Dueños"],
    ["GET/POST/PUT/DELETE", "/api/mascotas", "CRUD Mascotas"],
    ["GET/POST/PUT/DELETE", "/api/turnos", "CRUD Turnos"],
    ["GET/POST/PUT/DELETE", "/api/historiales", "CRUD Historial Clinico"],
    ["GET", "/api/dashboard/stats", "Estadisticas del dashboard"],
]
add_table(slide, Inches(1), Inches(1.5), Inches(11), Inches(4.5), headers, rows)

add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.8), "DELETE retorna 204 (sin contenido). Solo ADMIN puede eliminar registros.", 14, YELLOW, False, PP_ALIGN.LEFT)

# ===================== SLIDE 7: FRONTEND REACT =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "6. Frontend - React 19", 32, TEAL, True)

add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.3), DARK_CARD)
add_textbox(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.5), "Paginas Principales", 20, VIOLET, True)
add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5), Inches(4.2), [
    "- Login → auth dual login/registro",
    "- Dashboard → panel estadisticas",
    "- Duenos → CRUD completo (tabla + form)",
    "- Mascotas → CRUD completo",
    "- Turnos → CRUD con picker fecha",
    "- ClientPortal → vista cliente (tabs)",
], 15, LIGHT_GRAY)

add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), DARK_CARD)
add_textbox(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.5), "Tecnologias y Estado", 20, VIOLET, True)
add_bullet_list(slide, Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2), [
    "- React 19 (SPA con React Router v7)",
    "- Estado: useState + useEffect + localStorage",
    "- Tema oscuro: CSS Custom Properties",
    "- Colores: Teal (#14b8a6) + Violet (#8b5cf6)",
    "- Responsive: media query 768px",
    "- Deploy: Vercel (automatico desde GitHub)",
], 15, LIGHT_GRAY)

# ===================== SLIDE 8: AUTENTICACION Y SEGURIDAD =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "7. Autenticacion y Seguridad", 32, TEAL, True)

add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.3), DARK_CARD)
add_textbox(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.5), "Flujo JWT", 20, VIOLET, True)
add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5), Inches(4.2), [
    "1. Usuario ingresa credenciales",
    "2. POST /auth/login → backend verifica bcrypt",
    "3. Genera JWT: HS256, 60 min expira",
    "4. Payload: sub(usuario), rol, dueno_cedula",
    "5. Frontend guarda en localStorage",
    "6. Todas las peticiones incluyen:",
    "   Authorization: Bearer [token]",
], 15, LIGHT_GRAY)

add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), DARK_CARD)
add_textbox(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.5), "Roles y Permisos", 20, VIOLET, True)

headers = ["Operacion", "ADMIN", "USER"]
rows = [
    ["Crear Dueño/Mascota/Turno", "Si", "Si"],
    ["Ver Dashboard", "Si", "Si"],
    ["Ver Sus Propios Datos", "Si", "Si"],
    ["Eliminar registros", "Si", "No"],
    ["Crear/Eliminar Historial", "Si", "No"],
]
add_table(slide, Inches(7.1), Inches(2.4), Inches(5.2), Inches(3.5), headers, rows)

add_textbox(slide, Inches(7.1), Inches(6.1), Inches(5.2), Inches(0.5), "CORS: localhost:3000, :5173, vercel.app", 12, LIGHT_GRAY, False)

# ===================== SLIDE 9: AUTOMATIZACION =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "8. Automatizacion - Serenity BDD", 32, TEAL, True)

add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.3), DARK_CARD)
add_textbox(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.5), "Stack y Patron Screenplay", 20, VIOLET, True)
add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5), Inches(4.2), [
    "- Serenity BDD 5.3.9 + Cucumber 7.15",
    "- Selenium WebDriver + WebDriverManager",
    "- JUnit 4.13.2 como test runner",
    "",
    "Patron Screenplay:",
    "- Models: DuenoData, MascotaData, TurnoData",
    "- Tasks: 7 acciones (login, CRUD, etc.)",
    "- Questions: 5 validaciones",
    "- Pages: LoginPage, DashboardPage, etc.",
], 15, LIGHT_GRAY)

add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), DARK_CARD)
add_textbox(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.5), "8 Escenarios de Prueba", 20, VIOLET, True)

headers = ["#", "Escenario"]
rows = [
    ["1", "Registrar un dueno exitosamente"],
    ["2", "Registrar mascota con dueno existente"],
    ["3", "Registrar turno con mascota existente"],
    ["4", "Consultar mascotas registradas"],
    ["5", "Consultar turnos registrados"],
    ["6", "Dashboard muestra contadores"],
    ["7", "Editar dueno existente"],
    ["8", "Cerrar sesion del sistema"],
]
add_table(slide, Inches(7.1), Inches(2.4), Inches(5.2), Inches(4), headers, rows)

# ===================== SLIDE 10: BASE DE DATOS =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "9. Base de Datos - MySQL", 32, TEAL, True)

add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.3), DARK_CARD)
add_textbox(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.5), "Configuracion", 20, VIOLET, True)
add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5), Inches(4.2), [
    "- Motor: MySQL",
    "- Base de datos: veterinaria",
    "- Host: localhost:3306",
    "- Usuario: root",
    "- DDL Mode: update (auto)",
    "",
    "ORM: SQLAlchemy (mapea clases → tablas)",
    "Driver: pymysql",
], 15, LIGHT_GRAY)

add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), DARK_CARD)
add_textbox(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.5), "Relaciones y Restricciones", 20, VIOLET, True)
add_bullet_list(slide, Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2), [
    "PRIMARY KEY:",
    "  cedula (duenos), id (demas tablas)",
    "",
    "FOREIGN KEY:",
    "  mascotas → duenos",
    "  turnos → mascotas",
    "  historial → mascotas",
    "  usuarios → duenos",
    "",
    "UNIQUE: username en usuarios",
    "DEFAULT: rol = USER",
    "ON DELETE CASCADE en todas las FK",
], 14, LIGHT_GRAY)

# ===================== SLIDE 11: RESUMEN ESTADO =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "10. Resumen y Estado del Proyecto", 32, TEAL, True)

headers = ["Componente", "Estado", "Detalle"]
rows = [
    ["Backend FastAPI", "COMPLETO", "CRUD, Auth JWT, Dashboard"],
    ["Frontend React", "COMPLETO", "Login, CRUD, Portal, Dashboard"],
    ["Base de Datos MySQL", "COMPLETO", "5 tablas, relaciones OK"],
    ["Automatizacion E2E", "COMPLETO", "8 escenarios, reportes HTML"],
    ["App Movil", "PENDIENTE", "Carpeta vacia"],
    ["Documentacion", "PARCIAL", "Diagrama de clases UML"],
]
add_table(slide, Inches(1), Inches(1.5), Inches(11), Inches(4), headers, rows)

add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5), "Herramientas: Git/GitHub, VS Code, Postman, draw.io, MySQL Workbench, Vercel", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)

output_path = r'C:\Users\Michael\Documents\PetClinic\DOCUMENTACION\PetClinic_Resumen.pptx'
prs.save(output_path)
print(f'PowerPoint generado en: {output_path}')
